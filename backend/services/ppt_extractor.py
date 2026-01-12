import collections 
import pptx

def extract_text_from_ppt(ppt_file_path: str) -> str:
    """
    Extracts text from slides, including titles, body text, and speaker notes.
    Returns a single string with structured content.
    """
    prs = pptx.Presentation(ppt_file_path)
    
    text_content = []
    
    for i, slide in enumerate(prs.slides):
        slide_content = []
        
        # Slide Number
        slide_content.append(f"--- Slide {i+1} ---")
        
        # Title
        if slide.shapes.title and slide.shapes.title.text:
            slide_content.append(f"Title: {slide.shapes.title.text.strip()}")
            
        # Body Text
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                # Avoid duplicating title
                if shape == slide.shapes.title:
                    continue
                cleaned_text = shape.text.strip()
                if cleaned_text:
                    slide_content.append(cleaned_text)
        
        # Speaker Notes
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                slide_content.append(f"Notes: {notes}")
        
        text_content.append("\n".join(slide_content))
        
    return "\n\n".join(text_content)
